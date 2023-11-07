import fitz  # PyMuPDF
from PIL import Image
import io
from pptx import Presentation
from pptx.util import Inches
import cv2
import pytesseract
import numpy as np
import tempfile

# Function to check if an image contains text
def is_image_clear(image):
    text = pytesseract.image_to_string(image)
    return not text

# Input PDF file and output PPT file names
pdf_file = 'input.pdf'
ppt_file = 'output.pptx'

# Create a temporary directory to store image files
temp_dir = tempfile.TemporaryDirectory()

# Extract images from the PDF and filter out those with text
images = []
pdf_document = fitz.open(pdf_file)
for page_num in range(pdf_document.page_count):
    page = pdf_document.load_page(page_num)
    xref_list = page.get_images(full=True)
    for img_xref in xref_list:
        base_image = pdf_document.extract_image(img_xref[0])
        image_data = base_image["image"]
        image = Image.open(io.BytesIO(image_data))
        if is_image_clear(image):
            # Save the image to a temporary file
            temp_image_path = temp_dir.name + f'/image_{len(images)}.png'
            image.save(temp_image_path)
            images.append(temp_image_path)

# Create a PowerPoint presentation and add the filtered images
ppt = Presentation()
for image_path in images:
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Blank slide layout
    left = top = Inches(1)
    pic_width = Inches(8)
    pic_height = Inches(6)
    slide.shapes.add_picture(image_path, left, top, pic_width, pic_height)

# Save the PowerPoint presentation
ppt.save(ppt_file)
print(f'PowerPoint presentation saved as {ppt_file}')

# Clean up temporary files and directory
temp_dir.cleanup()
